import hashlib
import re
from pathlib import Path

from django.conf import settings
from django.core.management.base import BaseCommand
from django.utils.text import slugify

from apps.livia_assistant.models import LiviaKnowledgeItem


DEFAULT_CHUNK_SIZE = 3200
MIN_CHUNK_SIZE = 800


def normalize_extracted_text(text):
    text = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
    lines = [line for line in lines if line]
    return "\n".join(lines).strip()


def chunk_text(text, chunk_size=DEFAULT_CHUNK_SIZE, min_chunk_size=MIN_CHUNK_SIZE):
    cleaned = normalize_extracted_text(text)
    if not cleaned:
        return []

    paragraphs = [paragraph.strip() for paragraph in cleaned.split("\n") if paragraph.strip()]
    chunks = []
    current = ""
    for paragraph in paragraphs:
        candidate = f"{current}\n{paragraph}".strip() if current else paragraph
        if len(candidate) <= chunk_size:
            current = candidate
            continue
        if current:
            chunks.append(current)
            current = paragraph
        else:
            start = 0
            while start < len(paragraph):
                chunks.append(paragraph[start : start + chunk_size].strip())
                start += chunk_size
            current = ""
    if current:
        chunks.append(current)

    merged_chunks = []
    for chunk in chunks:
        if merged_chunks and len(chunk) < min_chunk_size:
            merged_chunks[-1] = f"{merged_chunks[-1]}\n{chunk}".strip()
            continue
        merged_chunks.append(chunk)
    return [chunk for chunk in merged_chunks if chunk]


def extract_pdf_text(path):
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("Dependência ausente para PDF. Instale com: pip install pypdf") from exc

    reader = PdfReader(str(path))
    pages = []
    for page in reader.pages:
        page_text = normalize_extracted_text(page.extract_text() or "")
        if page_text:
            pages.append(page_text)
    return "\n".join(pages)


def extract_pptx_text(path):
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Dependência ausente para PPTX. Instale com: pip install python-pptx") from exc

    presentation = Presentation(str(path))
    slides = []
    for slide in presentation.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                value = normalize_extracted_text(shape.text or "")
                if value:
                    parts.append(value)
        slide_text = normalize_extracted_text("\n".join(parts))
        if slide_text:
            slides.append(slide_text)
    return "\n".join(slides)


def build_keywords(file_path):
    stem = file_path.stem.replace("_", " ").replace("-", " ").strip().lower()
    stem_tokens = [token for token in re.findall(r"[a-z0-9]{3,}", stem) if token]
    base = ["xyron", "xyron robotics", "smart control brasil", "catalogo", "produto", "documentacao tecnica"]
    return " ".join(base + stem_tokens)


def build_slug_prefix(relative_path):
    stem_slug = slugify(relative_path.stem) or "documento"
    digest = hashlib.sha1(relative_path.as_posix().encode("utf-8")).hexdigest()[:8]
    prefix = f"xyron-doc-{stem_slug}-{digest}"
    return prefix[:190].rstrip("-")


class Command(BaseCommand):
    help = "Importa documentos locais PDF/PPTX da Xyron para a base da Lívia."

    def add_arguments(self, parser):
        parser.add_argument(
            "--path",
            default="",
            help="Caminho da pasta com arquivos Xyron. Padrão: docs/xyron/",
        )

    def handle(self, *args, **options):
        requested_path = (options.get("path") or "").strip()
        docs_path = Path(requested_path) if requested_path else Path(settings.BASE_DIR) / "docs" / "xyron"
        docs_path = docs_path.resolve()

        if not docs_path.exists() or not docs_path.is_dir():
            self.stdout.write(self.style.WARNING(f"Pasta não encontrada: {docs_path}"))
            return

        created_count = 0
        updated_count = 0
        ignored_count = 0
        error_count = 0

        supported_files = sorted(
            [*docs_path.rglob("*.pdf"), *docs_path.rglob("*.pptx")],
            key=lambda path: path.as_posix().lower(),
        )
        if not supported_files:
            self.stdout.write(self.style.WARNING(f"Nenhum arquivo .pdf/.pptx encontrado em: {docs_path}"))
            return

        for file_path in supported_files:
            relative_path = file_path.relative_to(docs_path)
            try:
                if file_path.suffix.lower() == ".pdf":
                    extracted = extract_pdf_text(file_path)
                else:
                    extracted = extract_pptx_text(file_path)
            except RuntimeError as exc:
                error_count += 1
                self.stderr.write(self.style.ERROR(f"[erro] {relative_path}: {exc}"))
                continue
            except Exception as exc:  # pragma: no cover - guard para arquivos corrompidos
                error_count += 1
                self.stderr.write(self.style.ERROR(f"[erro] {relative_path}: falha na extração ({exc})"))
                continue

            blocks = chunk_text(extracted)
            if not blocks:
                ignored_count += 1
                self.stdout.write(self.style.WARNING(f"[skip] {relative_path}: sem texto útil extraído"))
                continue

            slug_prefix = build_slug_prefix(relative_path)
            title_base = relative_path.stem.replace("_", " ").replace("-", " ").strip() or "Documento Xyron"
            keywords = build_keywords(relative_path)
            generated_slugs = []

            for index, block in enumerate(blocks, start=1):
                block_slug = f"{slug_prefix}-{index:03d}"[:200]
                generated_slugs.append(block_slug)
                title = f"{title_base} - Bloco {index}"[:180]
                content = f"Origem documento Xyron: {relative_path.as_posix()}\n\n{block}"

                _, created = LiviaKnowledgeItem.objects.update_or_create(
                    slug=block_slug,
                    defaults={
                        "title": title,
                        "category": LiviaKnowledgeItem.Category.TECHNICAL,
                        "content": content,
                        "keywords": keywords,
                        "priority": 82,
                        "is_active": True,
                    },
                )
                if created:
                    created_count += 1
                else:
                    updated_count += 1

            (
                LiviaKnowledgeItem.objects.filter(slug__startswith=f"{slug_prefix}-")
                .exclude(slug__in=generated_slugs)
                .update(is_active=False)
            )

        self.stdout.write(
            self.style.SUCCESS(
                "import_livia_xyron_docs concluído: "
                f"{created_count} criados, {updated_count} atualizados, "
                f"{ignored_count} ignorados, {error_count} erros."
            )
        )
